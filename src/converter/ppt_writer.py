# ppt_writer.py
import os
from typing import Optional
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(title: str = "Report", subtitle: str = "") -> Presentation:
    prs = Presentation()
    # Title slide layout is usually 0 (depends on template)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    # placeholder index 1 is often subtitle, but check existence
    subtitle_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title_placeholder.text = title
    if subtitle and subtitle_placeholder:
        subtitle_placeholder.text = subtitle
    return prs

def set_shape_text(shape, text: str, font_size: int = 18, bold: bool = False):
    tx = shape.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold

# --------- Slide creators ---------------------------------------------------

def row_to_text_slide(prs: Presentation, row: pd.Series, title_col: Optional[str] = None):
    """
    Create a text slide for a single row:
    - title: value from title_col (if provided) or the first non-empty field
    - body: bullet list of key: value for the other columns
    """
    layout = prs.slide_layouts[1]  # Title + Content layout (common)
    slide = prs.slides.add_slide(layout)

    # determine title text
    if title_col and title_col in row.index and not pd.isna(row[title_col]) and str(row[title_col]).strip() != "":
        title_text = str(row[title_col])
    else:
        # fallback: first non-empty column value or column name
        title_text = None
        for c in row.index:
            if not pd.isna(row[c]) and str(row[c]).strip() != "":
                title_text = str(row[c])
                break
        if title_text is None:
            # ultimate fallback: first column name
            title_text = str(row.index[0])

    # set slide title
    if slide.shapes.title:
        slide.shapes.title.text = title_text

    # get body placeholder (usually index 1)
    body_placeholder = None
    if len(slide.placeholders) > 1:
        body_placeholder = slide.placeholders[1]
    else:
        # try to find a content placeholder among shapes
        for shp in slide.shapes:
            if hasattr(shp, "text_frame"):
                body_placeholder = shp
                break

    if body_placeholder is None:
        # add textbox if no existing body placeholder
        body_placeholder = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))

    body = body_placeholder.text_frame
    body.clear()

    # add bullet lines
    for col in row.index:
        if title_col and col == title_col:
            continue
        val = row[col]
        if pd.isna(val) or str(val).strip() == "":
            continue
        p = body.add_paragraph()
        p.level = 0
        p.text = f"{col}: {val}"

def row_to_table_slide(prs: Presentation, row: pd.Series, title_col: Optional[str] = None, max_cols: int = 2):
    """
    Create a slide with a simple 2-column table: Field | Value
    """
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # Title (if present)
    if slide.shapes.title:
        if title_col and title_col in row.index and not pd.isna(row[title_col]) and str(row[title_col]).strip() != "":
            slide.shapes.title.text = str(row[title_col])
        else:
            for c in row.index:
                if not pd.isna(row[c]) and str(row[c]).strip() != "":
                    slide.shapes.title.text = str(row[c])
                    break

    # Prepare table data
    items = [(col, row[col]) for col in row.index if not pd.isna(row[col]) and str(row[col]).strip() != ""]
    if not items:
        # empty row; add a small textbox
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        set_shape_text(tx_box, "No data", font_size=20)
        return

    rows = len(items) + 1  # header + items
    cols = max_cols
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(0.6 + 0.2 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # header row
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"

    # fill table
    for i, (k, v) in enumerate(items, start=1):
        table.cell(i, 0).text = str(k)
        table.cell(i, 1).text = str(v)

    # Set font sizes for table cells
    for r in range(rows):
        for c in range(cols):
            for paragraph in table.cell(r, c).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

# --- Main export function --------------------------------------------------

def df_to_ppt(df: pd.DataFrame, out_path: str, title: str = "Auto Report", subtitle: str = "",
              title_col: Optional[str] = None, mode: str = "table", limit: Optional[int] = None):
    """
    Convert a DataFrame to PPT:
    - mode: 'table' or 'text' (per-row slide)
    - title_col: optional, column name used for slide title
    - limit: optional max number of rows to export (useful for testing)
    """
    if df is None or df.shape[0] == 0:
        prs = create_presentation(title, subtitle)
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "No data"
        prs.save(out_path)
        return out_path

    prs = create_presentation(title, subtitle)
    n = df.shape[0]
    if limit is not None:
        n = min(n, int(limit))

    for i in range(n):
        row = df.iloc[i]
        if mode == "text":
            row_to_text_slide(prs, row, title_col=title_col)
        else:
            row_to_table_slide(prs, row, title_col=title_col)

    # ensure output directory exists
    out_dir = os.path.dirname(out_path)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)

    prs.save(out_path)
    return out_path

# --- Example usage ---------------------------------------------------------

if __name__ == "__main__":
    demo_path = os.path.join("examples", "finance_sample.xlsx")
    df = pd.read_excel(demo_path, sheet_name=0)
    df = df.fillna("").astype(object)
    out = df_to_ppt(df, out_path=os.path.join("examples", "demo_presentation.pptx"),
                    title="Finance Sample", subtitle="Auto-generated", title_col="Asset", mode="table", limit=10)
    print("Saved:", out)
