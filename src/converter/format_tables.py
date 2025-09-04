import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Theme defaults (override via args if you want) ---
DEFAULT_HEADER_BG = RGBColor(0, 51, 102)
DEFAULT_HEADER_FG = RGBColor(255, 255, 255)
DEFAULT_ROW_EVEN_BG = RGBColor(235, 241, 222)
DEFAULT_ROW_ODD_BG  = RGBColor(255, 255, 255)

def _safe_layout(prs, preferred=5, fallback=1):
    if len(prs.slide_layouts) > preferred:
        return prs.slide_layouts[preferred]
    if len(prs.slide_layouts) > fallback:
        return prs.slide_layouts[fallback]
    return prs.slide_layouts[0]

def _auto_column_widths(df: pd.DataFrame, table, total_width_in=9.0, sample_rows=1000):
    lengths = []
    for col in df.columns:
        values = df[col].head(sample_rows)
        max_len = max([len(str(col))] + [len(str(v)) for v in values])
        lengths.append(max(1, max_len))
    total = sum(lengths) or 1
    for j, L in enumerate(lengths):
        frac = L / total
        table.columns[j].width = int(Inches(total_width_in) * frac)

def _set_cell_text(cell, text, bold=False, size=11, align=PP_ALIGN.LEFT, color=None):
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "" if text is None else str(text)
    run.font.bold = bold
    run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    p.alignment = align

def _apply_fill(cell, rgb: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb

def _add_light_borders(table):
    # python-pptx has limited border API via XML usually; do a light fallback:
    # We’ll set a minimal padding “look” using paragraph spacing only (portable & safe).
    # If you want true borders, you’d need an OXML hack; keeping this simple & robust.

    # Optional no-op placeholder to document intent.
    return

def _format_value(val, formatter):
    if formatter is None:
        return val
    try:
        return formatter(val)
    except Exception:
        return val

def add_dataframe_table_slide(
    prs: Presentation,
    df: pd.DataFrame,
    title: str = "Data",
    header_bg: RGBColor = DEFAULT_HEADER_BG,
    header_fg: RGBColor = DEFAULT_HEADER_FG,
    row_even_bg: RGBColor = DEFAULT_ROW_EVEN_BG,
    row_odd_bg: RGBColor = DEFAULT_ROW_ODD_BG,
    banded: bool = True,
    number_formatters: dict | None = None,  # e.g., {"Revenue": lambda x: f"{x:,.0f}"}
    max_rows: int | None = None,
    table_box=(0.5, 1.5, 9.0, 5.0),        # Inches: (left, top, width, height)
):
    """
    Adds a slide with a formatted table built from df and returns the Presentation.
    """
    # Clamp rows if requested
    work_df = df.copy()
    if max_rows is not None and max_rows >= 0:
        work_df = work_df.head(max_rows)

    rows, cols = work_df.shape
    slide_layout = _safe_layout(prs, preferred=5, fallback=1)
    slide = prs.slides.add_slide(slide_layout)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        # fallback textbox
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
        _set_cell_text(tx.text_frame.paragraphs[0]._parent, title, bold=True, size=24, align=PP_ALIGN.LEFT)

    left, top, width, height = map(Inches, table_box)
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Header
    for j, col_name in enumerate(work_df.columns):
        cell = table.cell(0, j)
        _set_cell_text(cell, col_name, bold=True, size=12, align=PP_ALIGN.CENTER, color=header_fg)
        _apply_fill(cell, header_bg)

    # Data
    from pandas.api.types import is_numeric_dtype
    fmt_map = number_formatters or {}

    for i in range(rows):
        row_bg = row_even_bg if (i % 2 == 0 and banded) else (row_odd_bg if banded else row_odd_bg)
        for j in range(cols):
            val = work_df.iat[i, j]
            col_name = work_df.columns[j]

            # Apply per-column formatter if given
            formatted = _format_value(val, fmt_map.get(col_name))

            # Alignment based on dtype
            align = PP_ALIGN.RIGHT if is_numeric_dtype(work_df.dtypes.iloc[j]) else PP_ALIGN.LEFT

            cell = table.cell(i + 1, j)
            _set_cell_text(cell, "" if pd.isna(formatted) else formatted, size=11, align=align)
            _apply_fill(cell, row_bg)

    # Column widths
    _auto_column_widths(work_df, table, total_width_in=table_box[2])

    # Optional “borders” placeholder
    _add_light_borders(table)

    return prs
